"""
Genera la presentación PowerPoint:
  RNN, LSTM y GRU: Pronóstico de Series de Tiempo
  Aprendizaje Profundo — UNSAM

Uso:
    python generar_presentacion.py
Salida:
    rnn_lstm_gru_presentacion.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ─── Paleta de colores ────────────────────────────────────────────────────────
C_BG       = RGBColor(0x1A, 0x1A, 0x2E)  # fondo oscuro (azul noche)
C_ACCENT   = RGBColor(0xE9, 0x4F, 0x37)  # rojo coral (acento)
C_ACCENT2  = RGBColor(0x39, 0xB5, 0xE0)  # azul claro
C_ACCENT3  = RGBColor(0x4C, 0xD9, 0x7A)  # verde
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xCC, 0xCC, 0xDD)
C_DARK_BOX = RGBColor(0x0F, 0x3A, 0x5F)  # caja azul oscura

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=20, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                           font_size=18, color=C_WHITE, bullet=False,
                           align=PP_ALIGN.LEFT):
    """lines: list of (text, bold, size_override_or_None, color_override_or_None)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, sz, col = item, False, None, None
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            sz   = item[2] if len(item) > 2 else None
            col  = item[3] if len(item) > 3 else None

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = align
        if bullet:
            para.level = 0

        run = para.add_run()
        prefix = "• " if bullet and text.strip() else ""
        run.text = prefix + text
        run.font.size = Pt(sz if sz else font_size)
        run.font.bold = bold
        run.font.color.rgb = col if col else color
        run.font.name = "Calibri"
    return txBox


def title_banner(slide, title, subtitle=None):
    """Top colored banner with title."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.5), C_DARK_BOX)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9),
                font_size=32, bold=True, color=C_ACCENT2,
                align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.5),
                    font_size=18, color=C_LIGHT, align=PP_ALIGN.LEFT)


def section_divider(prs, number, title):
    """Interstitial section divider slide."""
    slide = blank_slide(prs)
    fill_bg(slide, C_DARK_BOX)
    add_rect(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.08), C_ACCENT, )
    add_textbox(slide, f"Sección {number}", Inches(0.8), Inches(1.8),
                Inches(11), Inches(0.7), font_size=22, color=C_ACCENT, bold=False)
    add_textbox(slide, title, Inches(0.8), Inches(2.4),
                Inches(11.5), Inches(1.4), font_size=44, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)
    return slide


def code_box(slide, code_lines, left, top, width, height, font_size=13):
    """A dark box mimicking a code cell."""
    add_rect(slide, left, top, width, height, RGBColor(0x0D, 0x1B, 0x2A))
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code_lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = RGBColor(0xA8, 0xE6, 0xCF)


# ═════════════════════════════════════════════════════════════════════════════
# BUILD SLIDES
# ═════════════════════════════════════════════════════════════════════════════

prs = new_prs()

# ──────────────────────────────────────────────────────────────────────────────
# 1. PORTADA
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)

# Gradiente visual (simulated with rectangles)
add_rect(slide, Inches(0), Inches(0), Inches(4), SLIDE_H, RGBColor(0x0F, 0x3A, 0x5F))
add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, C_ACCENT)

add_textbox(slide, "Aprendizaje Profundo — UNSAM",
            Inches(4.3), Inches(1.2), Inches(8.5), Inches(0.6),
            font_size=18, color=C_LIGHT, italic=True)

add_textbox(slide, "RNN, LSTM y GRU",
            Inches(4.3), Inches(1.9), Inches(8.5), Inches(1.1),
            font_size=48, bold=True, color=C_WHITE)

add_textbox(slide, "Pronóstico de Series de Tiempo",
            Inches(4.3), Inches(2.95), Inches(8.5), Inches(0.7),
            font_size=28, color=C_ACCENT2)

add_rect(slide, Inches(4.3), Inches(3.7), Inches(7), Inches(0.05), C_ACCENT)

add_multiline_textbox(slide, [
    ("Dataset: Jena Climate (Max Planck Institute, 2009–2016)", False, 16, C_LIGHT),
    ("Framework: PyTorch  |  Idioma: Español", False, 16, C_LIGHT),
], Inches(4.3), Inches(3.85), Inches(8.5), Inches(0.9))

# Side decoration — architecture names
side_lines = [
    ("SimpleRNN", False, 20, C_ACCENT),
    ("", False, 8, None),
    ("LSTM", False, 20, C_ACCENT2),
    ("", False, 8, None),
    ("GRU", False, 20, C_ACCENT3),
]
add_multiline_textbox(slide, side_lines,
                      Inches(0.3), Inches(1.8), Inches(3.2), Inches(3.5),
                      align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# 2. AGENDA
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Agenda", "Estructura de la clase")

sections = [
    ("1", "Series de tiempo: conceptos y dataset de Jena"),
    ("2", "Preprocesamiento y ventanas deslizantes"),
    ("3", "Baselines (sentido común y MLP)"),
    ("4", "RNN: mecánica e implementación manual"),
    ("5", "LSTM: compuertas y estado de la celda"),
    ("6", "GRU: simplificación eficiente de la LSTM"),
    ("7", "Comparación final de modelos"),
]

col_w = Inches(5.8)
for idx, (num, title) in enumerate(sections):
    row = idx % 4
    col = idx // 4
    top  = Inches(1.7) + row * Inches(1.35)
    left = Inches(0.4) + col * Inches(6.5)

    add_rect(slide, left, top, Inches(0.5), Inches(0.5), C_ACCENT)
    add_textbox(slide, num, left, top + Inches(0.02), Inches(0.5), Inches(0.48),
                font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, left + Inches(0.6), top, col_w, Inches(0.55),
                font_size=18, color=C_LIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 1
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "1 — 2", "Series de tiempo y preprocesamiento")

# ──────────────────────────────────────────────────────────────────────────────
# 3. ¿QUÉ ES UNA SERIE DE TIEMPO?
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "¿Qué es una serie de tiempo?",
             "Mediciones ordenadas en el tiempo a intervalos regulares")

add_multiline_textbox(slide, [
    ("Ejemplos cotidianos:", True, 20, C_ACCENT2),
], Inches(0.4), Inches(1.65), Inches(5.5), Inches(0.5))

examples = [
    ("Temperatura cada 10 min", "Meteorología"),
    ("Precio de acción diario",  "Economía"),
    ("ECG (electrocardiograma)", "Medicina"),
    ("Ventas semanales",          "Logística"),
    ("Consumo eléctrico horario", "Ingeniería"),
]
for i, (var, domain) in enumerate(examples):
    top = Inches(2.1) + i * Inches(0.75)
    add_rect(slide, Inches(0.4), top + Inches(0.08), Inches(0.35), Inches(0.35), C_ACCENT)
    add_textbox(slide, var,    Inches(0.85), top, Inches(3.2), Inches(0.5),
                font_size=17, color=C_WHITE)
    add_textbox(slide, domain, Inches(4.1),  top, Inches(2.0), Inches(0.5),
                font_size=15, color=C_LIGHT, italic=True)

add_rect(slide, Inches(6.5), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

add_textbox(slide, "Tareas principales:", Inches(6.8), Inches(1.65),
            Inches(6.0), Inches(0.5), font_size=20, bold=True, color=C_ACCENT2)

tasks = [
    ("Pronóstico (Forecasting)",   "Predecir valores futuros a partir del pasado"),
    ("Detección de anomalías",     "Identificar comportamientos inusuales"),
    ("Clasificación",              "Asignar etiquetas a segmentos de la serie"),
    ("Detección de eventos",       "Encontrar ocurrencias de patrones específicos"),
]
for i, (task, desc) in enumerate(tasks):
    top = Inches(2.1) + i * Inches(1.1)
    add_rect(slide, Inches(6.8), top, Inches(5.9), Inches(0.95),
             RGBColor(0x0F, 0x3A, 0x5F))
    add_textbox(slide, task, Inches(7.0), top + Inches(0.05), Inches(5.5), Inches(0.4),
                font_size=17, bold=True, color=C_ACCENT3)
    add_textbox(slide, desc, Inches(7.0), top + Inches(0.45), Inches(5.5), Inches(0.4),
                font_size=14, color=C_LIGHT)

add_textbox(slide, "Este notebook → Pronóstico",
            Inches(6.8), Inches(6.4), Inches(6.0), Inches(0.5),
            font_size=17, bold=True, color=C_ACCENT)


# ──────────────────────────────────────────────────────────────────────────────
# 4. DATASET DE JENA
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Dataset climático de Jena",
             "Max Planck Institute for Biogeochemistry, Alemania (2009–2016)")

left_lines = [
    ("Características del dataset:", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("• 420.551 mediciones (cada 10 minutos)", False, 17, C_WHITE),
    ("• 14 variables climáticas", False, 17, C_WHITE),
    ("• 8 años de datos (2009–2016)", False, 17, C_WHITE),
    ("", False, 8, None),
    ("Variables incluidas:", True, 18, C_ACCENT2),
    ("", False, 8, None),
    ("• T (°C)     — temperatura", False, 16, C_LIGHT),
    ("• p (mbar)  — presión atmosférica", False, 16, C_LIGHT),
    ("• rh (%)    — humedad relativa", False, 16, C_LIGHT),
    ("• wv (m/s)  — velocidad del viento", False, 16, C_LIGHT),
    ("• wd (deg)  — dirección del viento", False, 16, C_LIGHT),
    ("• + 9 variables más...", False, 16, C_LIGHT),
]
add_multiline_textbox(slide, left_lines,
                      Inches(0.4), Inches(1.65), Inches(5.5), Inches(5.2))

add_rect(slide, Inches(6.5), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

right_lines = [
    ("Problema de pronóstico:", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("Dado el historial de las últimas 5 días", False, 17, C_WHITE),
    ("(120 horas, submuestreadas a 1/hora),", False, 17, C_WHITE),
    ("predecir la temperatura 24 horas después.", False, 17, C_WHITE),
    ("", False, 12, None),
    ("Partición temporal (respeta el orden):", True, 18, C_ACCENT2),
]
add_multiline_textbox(slide, right_lines,
                      Inches(6.8), Inches(1.65), Inches(6.0), Inches(3.5))

# Split table
splits = [("Entrenamiento", "50%", "~210.000"), ("Validación", "25%", "~105.000"), ("Test", "25%", "~105.000")]
colors_split = [C_ACCENT3, C_ACCENT2, C_ACCENT]
for i, (name, pct, n) in enumerate(splits):
    top = Inches(4.7) + i * Inches(0.62)
    add_rect(slide, Inches(6.8), top, Inches(2.2), Inches(0.55), colors_split[i])
    add_textbox(slide, name, Inches(6.85), top + Inches(0.07), Inches(2.1), Inches(0.4),
                font_size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_textbox(slide, pct,  Inches(9.1),  top + Inches(0.07), Inches(1.2), Inches(0.4),
                font_size=14, color=C_WHITE)
    add_textbox(slide, n,    Inches(10.4), top + Inches(0.07), Inches(2.0), Inches(0.4),
                font_size=13, color=C_LIGHT)

add_textbox(slide,
            "Importante: nunca usar val/test para normalizar.",
            Inches(6.8), Inches(6.6), Inches(6.0), Inches(0.5),
            font_size=14, italic=True, color=C_ACCENT)


# ──────────────────────────────────────────────────────────────────────────────
# 5. PREPROCESAMIENTO
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Preprocesamiento",
             "Normalización y ventanas deslizantes")

add_textbox(slide, "Normalización (solo con estadísticas del conjunto de entrenamiento):",
            Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.5),
            font_size=19, bold=True, color=C_ACCENT2)

code_box(slide,
         ["x̃ = (x - μ_train) / σ_train",
          "",
          "mean = raw_data[:num_train].mean(axis=0)",
          "std  = raw_data[:num_train].std(axis=0)",
          "data_norm = (raw_data - mean) / std"],
         Inches(0.4), Inches(2.1), Inches(12.5), Inches(1.55))

add_textbox(slide, "Ventanas deslizantes (JenaWindowDataset):",
            Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.5),
            font_size=19, bold=True, color=C_ACCENT2)

params = [
    ("sampling_rate = 6",    "1 muestra por hora  (dataset: 1 cada 10 min)"),
    ("sequence_length = 120", "5 días de historia  (120 horas)"),
    ("delay = 858",           "target: temperatura 24 h después del último timestep"),
]
for i, (param, desc) in enumerate(params):
    top = Inches(4.25) + i * Inches(0.65)
    add_rect(slide, Inches(0.4), top, Inches(3.0), Inches(0.52), C_DARK_BOX)
    add_textbox(slide, param, Inches(0.5), top + Inches(0.07),
                Inches(2.8), Inches(0.4), font_size=14,
                color=C_ACCENT3, bold=True)
    add_textbox(slide, "→  " + desc, Inches(3.55), top + Inches(0.1),
                Inches(9.5), Inches(0.4), font_size=15, color=C_LIGHT)

code_box(slide,
         ["# Cada muestra: X=(120, 14) → y=temperatura real en °C"],
         Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.65), font_size=13)


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 2
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "3", "Baselines: sentido común y MLP")

# ──────────────────────────────────────────────────────────────────────────────
# 6. BASELINE DE SENTIDO COMÚN
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Baseline de sentido común",
             "La cota que cualquier modelo debe superar")

add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.3), C_DARK_BOX)
add_textbox(slide, '"La temperatura en 24 horas será igual a la temperatura ahora."',
            Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.6),
            font_size=22, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Heurística sin aprendizaje — aprovecha la inercia de la temperatura",
            Inches(0.7), Inches(2.35), Inches(12.0), Inches(0.5),
            font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

add_textbox(slide, "Métrica: Error Absoluto Medio (MAE)",
            Inches(0.4), Inches(3.2), Inches(12.5), Inches(0.5),
            font_size=20, bold=True, color=C_ACCENT2)

add_rect(slide, Inches(3.0), Inches(3.7), Inches(7), Inches(0.9), C_DARK_BOX)
add_textbox(slide, "MAE  =  (1/n) × Σ |ŷᵢ - yᵢ|",
            Inches(3.0), Inches(3.72), Inches(7), Inches(0.9),
            font_size=26, color=C_ACCENT3, align=PP_ALIGN.CENTER)

add_multiline_textbox(slide, [
    ("Resultados del baseline:", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("Validación MAE ≈ 2.44 °C", False, 22, C_WHITE),
    ("Test       MAE ≈ 2.62 °C", False, 22, C_WHITE),
    ("", False, 8, None),
    ("¿Por qué es difícil de superar? La temperatura tiene mucha inercia", False, 16, C_LIGHT),
    ("diaria: el mañana se parece bastante al hoy.", False, 16, C_LIGHT),
], Inches(0.4), Inches(4.7), Inches(8.0), Inches(2.5))

add_rect(slide, Inches(9.0), Inches(4.7), Inches(4.0), Inches(1.8), C_ACCENT)
add_textbox(slide, "Meta: superar este valor con modelos de deep learning",
            Inches(9.1), Inches(4.8), Inches(3.8), Inches(1.5),
            font_size=17, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# 7. MLP BASELINE
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Baseline MLP (red densa)",
             "¿Qué pasa si simplemente aplanamos la secuencia?")

add_multiline_textbox(slide, [
    ("Arquitectura:", True, 20, C_ACCENT2),
    ("", False, 6, None),
    ("Input (batch, 120, 14)  →  Flatten  →  Dense(16, ReLU)  →  Dense(1)", False, 18, C_WHITE),
    ("", False, 12, None),
    ("Problema fundamental:", True, 20, C_ACCENT),
    ("", False, 6, None),
    ("• Al aplanar del tensor (batch, 120, 14) a (batch, 1680),", False, 17, C_LIGHT),
    ("  se destruye completamente el ORDER temporal.", False, 17, C_LIGHT),
    ("• El modelo no sabe cuál feature corresponde a 'hace 5 días'", False, 17, C_LIGHT),
    ("  versus 'hace 1 hora'.", False, 17, C_LIGHT),
    ("", False, 8, None),
    ("Resultado: MAE ~ 2.6–2.8 °C  (no supera el baseline)", False, 17, C_ACCENT),
], Inches(0.4), Inches(1.65), Inches(7.5), Inches(5.2))

add_rect(slide, Inches(8.0), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

code_box(slide, [
    "class MLPModel(nn.Module):",
    "    def __init__(self, seq_len=120, n_features=14):",
    "        super().__init__()",
    "        self.net = nn.Sequential(",
    "            nn.Flatten(),",
    "            nn.Linear(seq_len * n_features, 16),",
    "            nn.ReLU(),",
    "            nn.Linear(16, 1),",
    "        )",
    "",
    "    def forward(self, x):",
    "        return self.net(x)  # orden destruido",
], Inches(8.2), Inches(1.65), Inches(5.0), Inches(4.5))

add_textbox(slide,
            "Conclusión: necesitamos una arquitectura que RESPETE el orden temporal.",
            Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.65),
            font_size=17, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 3
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "4", "Redes Neuronales Recurrentes (RNN)")

# ──────────────────────────────────────────────────────────────────────────────
# 8. ¿QUÉ ES UNA RNN?
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Red Neuronal Recurrente (RNN)",
             "Una red con memoria: el estado oculto hₜ")

left_lines = [
    ("Idea clave:", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("Una RNN procesa una secuencia paso a", False, 17, C_WHITE),
    ("paso, manteniendo un estado oculto hₜ", False, 17, C_WHITE),
    ("que acumula información del pasado.", False, 17, C_WHITE),
    ("", False, 12, None),
    ("Ecuación fundamental:", True, 20, C_ACCENT2),
]
add_multiline_textbox(slide, left_lines, Inches(0.4), Inches(1.65), Inches(5.8), Inches(3.0))

add_rect(slide, Inches(0.4), Inches(4.15), Inches(5.8), Inches(0.85), C_DARK_BOX)
add_textbox(slide, "hₜ = tanh( Wₓ · xₜ + Wₕ · hₜ₋₁ + b )",
            Inches(0.5), Inches(4.2), Inches(5.6), Inches(0.75),
            font_size=20, bold=True, color=C_ACCENT3, align=PP_ALIGN.CENTER)

notation = [
    ("xₜ  — entrada en el paso t", False, 15, C_LIGHT),
    ("hₜ₋₁ — estado oculto anterior (memoria)", False, 15, C_LIGHT),
    ("Wₓ, Wₕ — matrices de pesos (aprendibles)", False, 15, C_LIGHT),
    ("hₜ  — nuevo estado (resume el pasado)", False, 15, C_LIGHT),
]
add_multiline_textbox(slide, notation, Inches(0.4), Inches(5.1), Inches(5.8), Inches(1.8))

add_rect(slide, Inches(6.5), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

code_box(slide, [
    "# Pseudocódigo del loop RNN",
    "h = zeros(hidden_size)      # estado inicial",
    "",
    "for x_t in secuencia:",
    "    h = tanh(W_x @ x_t",
    "             + W_h @ h + b)  # actualiza",
    "",
    "pred = W_out @ h             # usa solo h_T",
], Inches(6.7), Inches(1.65), Inches(6.3), Inches(3.5))

add_multiline_textbox(slide, [
    ("Ventaja vs MLP:", True, 18, C_ACCENT2),
    ("• Preserva el orden temporal", False, 16, C_WHITE),
    ("• Puede procesar secuencias de longitud variable", False, 16, C_WHITE),
    ("• Los parámetros se comparten en todos los pasos", False, 16, C_WHITE),
], Inches(6.7), Inches(5.2), Inches(6.3), Inches(2.0))


# ──────────────────────────────────────────────────────────────────────────────
# 9. RNN DESDE CERO (código)
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "RNN manual en PyTorch",
             "Loop recurrente explícito — exactamente lo que hace nn.RNN internamente")

code_box(slide, [
    "class ManualRNNModel(nn.Module):",
    "    def __init__(self, input_size, hidden_size):",
    "        super().__init__()",
    "        self.W_x = nn.Linear(input_size,  hidden_size, bias=False)",
    "        self.W_h = nn.Linear(hidden_size, hidden_size, bias=True)",
    "        self.fc  = nn.Linear(hidden_size, 1)",
    "",
    "    def forward(self, x):              # x: (batch, seq_len, features)",
    "        h = torch.zeros(batch, hidden) # estado inicial",
    "",
    "        for t in range(seq_len):       # ← el loop recurrente",
    "            x_t = x[:, t, :]",
    "            h = torch.tanh(            # ecuación de la RNN",
    "                    self.W_x(x_t) + self.W_h(h)",
    "                )",
    "",
    "        return self.fc(h)              # pronóstico a partir de h_T",
], Inches(0.4), Inches(1.65), Inches(7.5), Inches(5.5))

add_rect(slide, Inches(8.1), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

add_multiline_textbox(slide, [
    ("Puntos clave:", True, 19, C_ACCENT2),
    ("", False, 6, None),
    ("• El loop for t es la RNN.", False, 17, C_WHITE),
    ("• h se actualiza en CADA paso.", False, 17, C_WHITE),
    ("• Solo usamos el estado final h_T", False, 17, C_WHITE),
    ("  para hacer el pronóstico.", False, 17, C_WHITE),
    ("", False, 12, None),
    ("Versión de PyTorch (equivalente):", True, 19, C_ACCENT2),
], Inches(8.3), Inches(1.65), Inches(4.8), Inches(3.5))

code_box(slide, [
    "self.rnn = nn.RNN(input_size,",
    "                  hidden_size,",
    "                  batch_first=True)",
    "_, h_n = self.rnn(x)",
    "return self.fc(h_n.squeeze(0))",
], Inches(8.3), Inches(4.9), Inches(4.8), Inches(1.8))


# ──────────────────────────────────────────────────────────────────────────────
# 10. PROBLEMA DE LOS GRADIENTES
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "El problema de los gradientes que desaparecen",
             "¿Por qué la RNN simple falla con secuencias largas?")

add_textbox(slide, "Backpropagation Through Time (BPTT):",
            Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.5),
            font_size=20, bold=True, color=C_ACCENT2)

add_rect(slide, Inches(0.4), Inches(2.15), Inches(12.5), Inches(1.0), C_DARK_BOX)
add_textbox(slide, "∂L/∂h₀  =  ∂L/∂hₜ · ∏ᵢ ∂hᵢ/∂hᵢ₋₁  ≈  ∂L/∂hₜ · (Wₕ)ᵀ",
            Inches(0.6), Inches(2.2), Inches(12.0), Inches(0.9),
            font_size=22, color=C_ACCENT3, align=PP_ALIGN.CENTER)

add_multiline_textbox(slide, [
    ("Consecuencia:", True, 20, C_ACCENT),
    ("", False, 6, None),
    ("• Si los valores propios de Wₕ < 1:  gradiente → 0  (desaparece)", False, 18, C_WHITE),
    ("• Para T=120 pasos: (Wₕ)¹²⁰ ≈ 0,  los pasos lejanos no aprenden nada.", False, 18, C_WHITE),
    ("", False, 12, None),
    ("• Si los valores propios de Wₕ > 1:  gradiente → ∞ (explota)", False, 18, C_WHITE),
    ("  → clipping de gradiente es una solución parcial.", False, 16, C_LIGHT),
], Inches(0.4), Inches(3.2), Inches(12.5), Inches(2.8))

add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.85), C_DARK_BOX)
add_textbox(slide, "Solución: diseñar arquitecturas con 'autopistas' para el gradiente → LSTM y GRU",
            Inches(0.5), Inches(6.15), Inches(12.2), Inches(0.75),
            font_size=19, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 4
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "5", "LSTM — Long Short-Term Memory")

# ──────────────────────────────────────────────────────────────────────────────
# 11. LSTM INTUICIÓN
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "LSTM: intuición",
             "Hochreiter & Schmidhuber (1997) — solución al gradiente que desaparece")

add_rect(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.15), C_DARK_BOX)
add_textbox(slide,
            "Imaginá una cinta transportadora (estado de la celda cₜ) que corre en "
            "paralelo a la secuencia.",
            Inches(0.6), Inches(1.7), Inches(12.0), Inches(0.5),
            font_size=18, italic=True, color=C_LIGHT)
add_textbox(slide,
            "La información puede 'subirse', 'bajarse' o continuar intacta. Las compuertas controlan el flujo.",
            Inches(0.6), Inches(2.2), Inches(12.0), Inches(0.5),
            font_size=18, italic=True, color=C_LIGHT)

gates = [
    ("fₜ", "Forget",    C_ACCENT,  "¿Cuánto del estado anterior cₜ₋₁ conservar?"),
    ("iₜ", "Input",     C_ACCENT2, "¿Cuánta información nueva agregar?"),
    ("c̃ₜ", "Candidate", C_ACCENT3, "Propuesta de nueva información (tanh)"),
    ("oₜ", "Output",    RGBColor(0xFF, 0xD7, 0x00), "¿Cuánto del estado exponer como hₜ?"),
]

for i, (sym, name, col, desc) in enumerate(gates):
    top  = Inches(3.0) + i * Inches(1.0)
    add_rect(slide, Inches(0.4), top, Inches(0.55), Inches(0.75), col)
    add_textbox(slide, sym, Inches(0.4), top + Inches(0.12), Inches(0.55), Inches(0.55),
                font_size=22, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_textbox(slide, name, Inches(1.1), top + Inches(0.02), Inches(2.0), Inches(0.4),
                font_size=18, bold=True, color=col)
    add_textbox(slide, desc, Inches(1.1), top + Inches(0.42), Inches(11.5), Inches(0.4),
                font_size=16, color=C_LIGHT)

add_rect(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.05), C_ACCENT)
add_textbox(slide, "Ecuación del estado de la celda:   cₜ  =  fₜ ⊙ cₜ₋₁  +  iₜ ⊙ c̃ₜ",
            Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
            font_size=18, bold=True, color=C_ACCENT3, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# 12. LSTM ECUACIONES
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "LSTM: ecuaciones completas",
             "4 proyecciones lineales + activaciones no lineales")

eqs = [
    ("fₜ = σ( Wxf·xₜ + Whf·hₜ₋₁ + bf )", "forget gate — olvido del pasado",   C_ACCENT),
    ("iₜ = σ( Wxi·xₜ + Whi·hₜ₋₁ + bi )", "input gate  — qué nueva info entra", C_ACCENT2),
    ("c̃ₜ = tanh( Wxg·xₜ + Whg·hₜ₋₁ + bg )", "candidate  — propuesta nueva",     C_ACCENT3),
    ("oₜ = σ( Wxo·xₜ + Who·hₜ₋₁ + bo )", "output gate — qué se expone",        RGBColor(0xFF, 0xD7, 0x00)),
    ("cₜ = fₜ ⊙ cₜ₋₁  +  iₜ ⊙ c̃ₜ",       "nueva celda  — ← autopista del grad",C_WHITE),
    ("hₜ = oₜ ⊙ tanh( cₜ )",              "nuevo estado oculto",                C_WHITE),
]

for i, (eq, label, col) in enumerate(eqs):
    top = Inches(1.65) + i * Inches(0.85)
    add_rect(slide, Inches(0.4), top, Inches(7.0), Inches(0.72), C_DARK_BOX)
    add_textbox(slide, eq, Inches(0.55), top + Inches(0.1),
                Inches(6.7), Inches(0.55), font_size=18, bold=True, color=col)
    add_textbox(slide, "→  " + label, Inches(7.55), top + Inches(0.18),
                Inches(5.4), Inches(0.4), font_size=15, color=C_LIGHT)

add_rect(slide, Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.05), C_ACCENT2)
add_textbox(slide,
            "Clave: cₜ = fₜ ⊙ cₜ₋₁ + iₜ ⊙ c̃ₜ  — si fₜ≈1 e iₜ≈0 → gradiente fluye sin atenuarse",
            Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.45),
            font_size=15, italic=True, color=C_ACCENT3, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# 13. LSTM CÓDIGO
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "LSTM manual en PyTorch",
             "Compuertas explícitas — exactamente lo que hace nn.LSTM")

code_box(slide, [
    "for t in range(seq_len):",
    "    x_t   = x[:, t, :]",
    "    gates = W_x(x_t) + W_h(h)     # proyección conjunta 4H",
    "",
    "    g_f, g_i, g_g, g_o = gates.chunk(4, dim=-1)",
    "",
    "    f = torch.sigmoid(g_f)         # forget gate",
    "    i = torch.sigmoid(g_i)         # input gate",
    "    g = torch.tanh(g_g)            # candidate",
    "    o = torch.sigmoid(g_o)         # output gate",
    "",
    "    c = f * c  +  i * g            # actualiza la 'cinta'",
    "    h = o * torch.tanh(c)          # nuevo estado oculto",
], Inches(0.4), Inches(1.65), Inches(7.2), Inches(5.2))

add_rect(slide, Inches(7.7), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

add_multiline_textbox(slide, [
    ("Observación pedagógica:", True, 18, C_ACCENT2),
    ("", False, 8, None),
    ("La proyección se hace de tamaño", False, 16, C_WHITE),
    ("4*hidden_size y luego se divide", False, 16, C_WHITE),
    ("en 4 partes (chunk).", False, 16, C_WHITE),
    ("Esto es más eficiente que 4", False, 16, C_WHITE),
    ("multiplicaciones separadas.", False, 16, C_WHITE),
    ("", False, 12, None),
    ("Versión con nn.LSTM:", True, 18, C_ACCENT2),
], Inches(7.9), Inches(1.65), Inches(5.0), Inches(3.5))

code_box(slide, [
    "lstm = nn.LSTM(input_size,",
    "               hidden_size,",
    "               batch_first=True)",
    "_, (h_n, c_n) = lstm(x)",
    "return fc(h_n.squeeze(0))",
], Inches(7.9), Inches(5.0), Inches(5.0), Inches(1.8))


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 5
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "6", "GRU — Gated Recurrent Unit")

# ──────────────────────────────────────────────────────────────────────────────
# 14. GRU
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "GRU — Gated Recurrent Unit",
             "Cho et al. (2014) — simplificación eficiente de la LSTM")

add_multiline_textbox(slide, [
    ("Idea: fusionar forget e input gates en una sola (update gate zₜ).", False, 17, C_WHITE),
    ("Eliminar el estado de celda cₜ separado → solo un hₜ.", False, 17, C_WHITE),
    ("Resultado: ⅓ menos parámetros que LSTM, rendimiento similar.", False, 17, C_ACCENT3),
], Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.3))

eqs_gru = [
    ("zₜ = σ( Wxz·xₜ + Whz·hₜ₋₁ + bz )", "update gate — ¿cuánto actualizar?", C_ACCENT2),
    ("rₜ = σ( Wxr·xₜ + Whr·hₜ₋₁ + br )", "reset gate  — ¿cuánto del pasado usar?", C_ACCENT3),
    ("h̃ₜ = tanh( Wxh·xₜ + Whh·(rₜ⊙hₜ₋₁) )", "candidate   — nueva información propuesta", C_WHITE),
    ("hₜ = (1−zₜ) ⊙ hₜ₋₁  +  zₜ ⊙ h̃ₜ", "nuevo estado — mezcla pasado y nuevo", C_ACCENT),
]

for i, (eq, label, col) in enumerate(eqs_gru):
    top = Inches(2.9) + i * Inches(0.85)
    add_rect(slide, Inches(0.4), top, Inches(7.0), Inches(0.72), C_DARK_BOX)
    add_textbox(slide, eq, Inches(0.55), top + Inches(0.1),
                Inches(6.7), Inches(0.55), font_size=17, bold=True, color=col)
    add_textbox(slide, "→  " + label, Inches(7.55), top + Inches(0.18),
                Inches(5.4), Inches(0.4), font_size=14, color=C_LIGHT)

add_textbox(slide, "Intuición de hₜ = (1−zₜ)·hₜ₋₁ + zₜ·h̃ₜ:",
            Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.5),
            font_size=17, bold=True, color=C_ACCENT2)
add_textbox(slide,
            "Si zₜ ≈ 0  →  hₜ ≈ hₜ₋₁  (recuerda el pasado)   |   Si zₜ ≈ 1  →  hₜ ≈ h̃ₜ  (actualiza con info nueva)",
            Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.45),
            font_size=16, italic=True, color=C_WHITE)


# ──────────────────────────────────────────────────────────────────────────────
# 15. COMPARACIÓN LSTM vs GRU
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Comparación: RNN vs LSTM vs GRU",
             "¿Cuándo usar cada arquitectura?")

# Header row
headers = ["",  "SimpleRNN", "LSTM", "GRU"]
hcolors = [C_BG, C_ACCENT, C_ACCENT2, C_ACCENT3]
col_starts = [Inches(0.4), Inches(2.5), Inches(5.65), Inches(9.3)]
col_widths  = [Inches(2.0), Inches(3.0),Inches(3.5), Inches(3.5)]

for j, (h, hc, x, w) in enumerate(zip(headers, hcolors, col_starts, col_widths)):
    if j > 0:
        add_rect(slide, x, Inches(1.65), w, Inches(0.6), hc)
    add_textbox(slide, h, x, Inches(1.65), w, Inches(0.6),
                font_size=18, bold=True, color=C_BG if j > 0 else C_WHITE,
                align=PP_ALIGN.CENTER)

rows = [
    ("Estado",      "hₜ",          "hₜ , cₜ",           "hₜ"),
    ("Compuertas",  "—",           "3 (f, i, o)",        "2 (z, r)"),
    ("Parámetros",  "2(d+u)u",     "4(d+u)u",            "3(d+u)u"),
    ("Gradientes",  "Desaparecen", "Fluyen por cₜ",      "Fluyen por hₜ"),
    ("Velocidad",   "Rápida",      "Más lenta",           "Intermedia"),
    ("Uso ideal",   "Problemas\nsimples", "Seq. largas\ny complejas", "Balance\ncapacidad-velocidad"),
]

row_colors = [C_DARK_BOX, RGBColor(0x12, 0x2A, 0x45),
              C_DARK_BOX, RGBColor(0x12, 0x2A, 0x45),
              C_DARK_BOX, RGBColor(0x12, 0x2A, 0x45)]

for i, (label, rnn_v, lstm_v, gru_v) in enumerate(rows):
    top = Inches(2.3) + i * Inches(0.78)
    vals = [label, rnn_v, lstm_v, gru_v]
    for j, (val, x, w) in enumerate(zip(vals, col_starts, col_widths)):
        add_rect(slide, x, top, w, Inches(0.75), row_colors[i])
        add_textbox(slide, val, x + Inches(0.05), top + Inches(0.07),
                    w - Inches(0.1), Inches(0.65),
                    font_size=14, bold=(j == 0),
                    color=C_LIGHT if j == 0 else C_WHITE,
                    align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# Sección divider 6
# ──────────────────────────────────────────────────────────────────────────────
section_divider(prs, "7", "Resultados y comparación de modelos")

# ──────────────────────────────────────────────────────────────────────────────
# 16. RESULTADOS
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Resultados: comparación de modelos",
             "Dataset Jena — MAE en grados Celsius (°C), FAST_MODE = False")

models_res = [
    ("Baseline sentido común",  "~2.44", "~2.62", "—",         C_LIGHT),
    ("MLP (red densa)",         "~2.65", "~2.75", "~27.000",   C_LIGHT),
    ("RNN Manual",              "~2.60", "~2.70", "~800",      C_ACCENT),
    ("nn.RNN",                  "~2.58", "~2.68", "~800",      C_ACCENT),
    ("LSTM Manual",             "~2.50", "~2.60", "~3.200",    C_ACCENT2),
    ("nn.LSTM + dropout",       "~2.45", "~2.55", "~3.200",    C_ACCENT2),
    ("GRU Manual",              "~2.48", "~2.58", "~2.400",    C_ACCENT3),
    ("GRU Apilada + dropout",   "~2.42", "~2.52", "~4.800",    C_ACCENT3),
]

# Header
header_cells = ["Modelo", "Val MAE", "Test MAE", "Parámetros"]
hx = [Inches(0.4), Inches(5.6), Inches(7.8), Inches(10.2)]
hw = [Inches(5.1), Inches(2.1), Inches(2.3), Inches(2.8)]
for h, x, w in zip(header_cells, hx, hw):
    add_rect(slide, x, Inches(1.65), w, Inches(0.55), C_DARK_BOX)
    add_textbox(slide, h, x + Inches(0.05), Inches(1.67), w, Inches(0.5),
                font_size=16, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

for i, (name, val, test, params, col) in enumerate(models_res):
    top = Inches(2.25) + i * Inches(0.6)
    bg = C_DARK_BOX if i % 2 == 0 else RGBColor(0x12, 0x2A, 0x45)
    add_rect(slide, Inches(0.4), top, Inches(5.1), Inches(0.57), bg)
    add_rect(slide, Inches(5.6), top, Inches(2.1), Inches(0.57), bg)
    add_rect(slide, Inches(7.8), top, Inches(2.3), Inches(0.57), bg)
    add_rect(slide, Inches(10.2), top, Inches(2.8), Inches(0.57), bg)

    add_textbox(slide, name, Inches(0.5), top + Inches(0.08), Inches(4.9), Inches(0.42),
                font_size=14, color=col)
    add_textbox(slide, val,  Inches(5.6),  top + Inches(0.08), Inches(2.0), Inches(0.42),
                font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, test, Inches(7.8),  top + Inches(0.08), Inches(2.2), Inches(0.42),
                font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, params, Inches(10.2), top + Inches(0.08), Inches(2.7), Inches(0.42),
                font_size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# 17. CONCLUSIONES
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Conclusiones y próximos pasos")

left_cols = [
    ("Lecciones clave:", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("1. Establecer un baseline primero.", True, 17, C_WHITE),
    ("   La heurística simple puede ser difícil de superar.", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("2. Las RNNs preservan el orden temporal.", True, 17, C_WHITE),
    ("   A diferencia del MLP que aplana y destruye la secuencia.", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("3. Las compuertas son la clave.", True, 17, C_WHITE),
    ("   fₜ⊙cₜ₋₁ (LSTM) y (1-zₜ)⊙hₜ₋₁ (GRU) crean autopistas", False, 15, C_LIGHT),
    ("   para el gradiente → resuelven el vanishing gradient.", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("4. Implementar desde cero vale la pena.", True, 17, C_WHITE),
    ("   Ver el loop y cada compuerta explica qué hace nn.LSTM.", False, 15, C_LIGHT),
]
add_multiline_textbox(slide, left_cols, Inches(0.4), Inches(1.65), Inches(6.5), Inches(5.4))

add_rect(slide, Inches(7.1), Inches(1.6), Inches(0.05), Inches(5.5), C_ACCENT)

right_cols = [
    ("¿Cómo continuar?", True, 20, C_ACCENT2),
    ("", False, 8, None),
    ("• Bidirectional RNNs:", True, 17, C_WHITE),
    ("  Procesar en ambas direcciones (NLP)", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("• Attention y Transformers:", True, 17, C_WHITE),
    ("  La arquitectura que reemplazó a las RNNs", False, 15, C_LIGHT),
    ("  en la mayoría de los tareas de NLP.", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("• Regularización avanzada:", True, 17, C_WHITE),
    ("  recurrent_dropout, weight_decay", False, 15, C_LIGHT),
    ("", False, 6, None),
    ("• Feature engineering:", True, 17, C_WHITE),
    ("  Hora del día, mes del año, etc.", False, 15, C_LIGHT),
]
add_multiline_textbox(slide, right_cols, Inches(7.3), Inches(1.65), Inches(5.7), Inches(5.4))


# ──────────────────────────────────────────────────────────────────────────────
# 18. EJERCICIOS
# ──────────────────────────────────────────────────────────────────────────────
slide = blank_slide(prs)
fill_bg(slide, C_BG)
title_banner(slide, "Ejercicios propuestos",
             "Para practicar en el notebook 02_rnn_lstm_gru_timeseries.ipynb")

exercises = [
    ("1", "Poné FAST_MODE = False y entrenás más épocas.\n¿Cuánto mejoran los modelos?"),
    ("2", "Aumentá HIDDEN a 64 o 128 en la LSTM.\n¿Cuándo empieza a overfittear?"),
    ("3", "Implementá un ManualLSTMModel bidireccional procesando la secuencia\nen orden inverso y promediando con el modelo normal."),
    ("4", "Probá sequence_length = 240 (10 días de historia).\n¿Mejora el pronóstico?"),
    ("5", "Modificá ManualLSTMModel para devolver los estados en cada timestep\ny visualizarlos como heatmap (igual que la RNN en rnn_example.ipynb)."),
]

for i, (num, text) in enumerate(exercises):
    top = Inches(1.65) + i * Inches(1.1)
    add_rect(slide, Inches(0.4), top, Inches(0.65), Inches(0.65), C_ACCENT)
    add_textbox(slide, num, Inches(0.4), top + Inches(0.05), Inches(0.65), Inches(0.6),
                font_size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, text, Inches(1.2), top + Inches(0.05), Inches(11.7), Inches(0.9),
                font_size=16, color=C_LIGHT)


# ──────────────────────────────────────────────────────────────────────────────
# GUARDAR
# ──────────────────────────────────────────────────────────────────────────────
output_path = "rnn_lstm_gru_presentacion.pptx"
prs.save(output_path)
print(f"Presentación guardada en: {output_path}")
print(f"Total de slides: {len(prs.slides)}")
